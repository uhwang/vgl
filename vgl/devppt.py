'''
    devpptx.py
    
    7/23/2023
    
    Line
    
    Use Case 1 - Specify all points in slide coordinates
    A very simple case can be had by specifying the verticies in slide coordinates:
    
    >>> freeform_builder = shapes.build_freeform(Inches(1.5), Inches(1))
    >>> freeform_builder.add_line_segments((
    ...     (Inches(2),   Inches(2)),
    ...     (Inches(1),   Inches(2)),
    ...     (Inches(1.5), Inches(1)),
    ... ))
    >>> freeform_shape = freeform_builder.convert_to_shape()
    >>> freeform_shape.left.inches, freeform_shape.top.inches
    1.0, 1.0
    >>> freeform_shape.width.inches, freeform_shape.height.inches
    1.0, 1.0
    
    Polygon
    
    >>> vertices = ((100, 200), (200, 300), (300, 400))
    >>> freeform_builder.add_line_segments(vertices, close=True)

    
    chat.openai.com
    ================
    
    line_shape = slide.shapes.add_shape(9, left, top, width, height)
    line_shape.shadow.inherit = False
    line_shape.shadow.blur_radius = 0
    line_shape.shadow.distance = 0
    
    after creating the line shape using slide.shapes.add_shape(), we access the shadow property of the shape and set the inherit, blur_radius, and distance attributes to disable the shadow effect. By setting inherit to False, we make sure that the shape's shadow properties are independent and not inherited from any style.
    
    With these changes, the line will be drawn without any shadow effect in the resulting PowerPoint presentation.
    
'''
import collections 
import collections.abc
import pptx
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

from . import color
from . import device
from . import drvemf 
from . import linepat
from . import patline
from . import gdiobj

_SHAPE_LINE = 9

def Polyline(slide, x, y, lcol, lthk, lpat, fcol, closed):

    free_form = slide.shapes.build_freeform(Inches(x[0]), Inches(y[0]))
    free_form.add_line_segments(
        [(Inches(x1), Inches(y1)) for x1, y1 in zip(x[1:],y[1:])],
        close=closed
    )
    
    line_shape = free_form.convert_to_shape()
    
    if closed == False: 
        line_shape.fill.background()
    elif fcol:
        line_shape.fill.solid()
        line_shape.fill.back_color = RGBColor(fcol.r, fcol.g, fcol.b)
    
    if lcol:
        if fcol and lcol == fcol:
            line_shape.line.color.rgb = RGBColor(fcol.r, fcol.g, fcol.b)
            line_shape.line.width = Pt(lthk)
        else:
            line_shape.line.color.rgb = RGBColor(lcol.r, lcol.g, lcol.b)
            line_shape.line.width = Pt(lthk)
        
    line_shape.shadow.inherit     = False
    line_shape.shadow.blur_radius = 0
    line_shape.shadow.distance    = 0

    
# Not working. 

#def Line(slide, x1, y1, x2, y2, lcol, lthk, lpat):
#    line_shape= slide.shapes.add_shape(_SHAPE_LINE, 
#                                        Inches(x1), 
#                                        Inches(y1), 
#                                        Inches(x2-x1), 
#                                        Inches(y2-y1))
#    line_shape.line.color.rgb = RGBColor(lcol.r, lcol.g, lcol.b)
#    line_shape.line.width = Pt(lthk)
#    line_shape.fill.background()
#    line_shape.shadow.inherit = False
#    line_shape.shadow.blur_radius = 0
#    line_shape.shadow.distance = 0
#    line_shape.rotation = -45000  # Rotation angle is in 1/60000 of a degree
    
    
class DevicePPT(device.DeviceVector):
    def __init__(self, fname, gbox):
        self.fname = fname
        self.ppt = pptx.Presentation()
        self.blank_slide_layout = self.ppt.slide_layouts[6]
        self.slide = self.ppt.slides.add_slide(self.blank_slide_layout)
        self.pen = gdiobj.Pen()
        self.brush = gdiobj.Brush()
        
    def set_device(self, frm, extend=device._FIT_NONE):
        self.frm = frm
        self.set_plot(frm,extend)
        
    def close(self):
        self.ppt.save(self.fname)
            
    def make_pen(self, color, thk):
        self.pen.lcol = color
        self.pen.lthk = thk
    
    def delete_pen(self):
        self.pen.lcol = None
        self.pen.lthk = None
        
    def line(self, sx, sy, ex, ey, lcol=None, lthk=None, lpat=linepat._PAT_SOLID, viewport=False):
        if self.pen:
            lcol = self.pen.lcol
            lthk = self.pen.lthk
        self.polyline((sx,ex), (sy,ey), lcol, lthk, lpat, viewport=False)
        
    def lline(self, sx, sy, ex, ey, lcol=None, lthk=None, lpat=linepat._PAT_SOLID, viewport=True):
        if self.pen:
            lcol = self.pen.lcol
            lthk = self.pen.lthk
        self.polyline((sx,ex), (sy,ey), lcol, lthk, lpat, viewport=True)
        
    def polyline(self, x, y, lcol=None, lthk=None, lpat=linepat._PAT_SOLID, closed=False, viewport=False):
        pat_inst = isinstance(lpat, linepat.LinePattern)

        if lthk: _lthk = lthk*self.frm.hgt()
        else: lthk = 0
        
        if pat_inst == False and lcol:
            if viewport:
                px, py = x, y
            else:
                px = [self._x_viewport(xx) for xx in x]
                py = [self._y_viewport(yy) for yy in y]
            Polyline(self.slide, px, py, lcol, lthk, lpat, None, closed)
    
        if lcol and pat_inst:
            xp, yp = x, y
                
            if viewport:
                pat_seg = patline.get_pattern_line(self, xp, yp, lpat.pat_len, lpat.pat_t, viewport=True)
            else:
                pat_seg = patline.get_pattern_line(self, xp, yp, lpat.pat_len, lpat.pat_t)

            for p1 in pat_seg:
                x2 = [p2[0] for p2 in p1 ]
                y2 = [p2[1] for p2 in p1 ]
                Polyline(self.slide, x2, y2, lcol, lthk, lpat, None, closed=False)
                
    def lpolyline(self, x, y, lcol=None, lthk=None, lpat=linepat._PAT_SOLID, closed=False, viewport=False):
        self.polyline(x,y,lcol,lthk,lpat,closed,viewport=True)
    
    def polygon(self, x, y, lcol=None, lthk=None, fcol=None, lpat=linepat._PAT_SOLID, viewport=False):
        
        pat_inst = isinstance(lpat, linepat.LinePattern)

        if lthk: _lthk = lthk*self.frm.hgt()
        else: lthk = 0
        
        if pat_inst == False:
            if viewport:
                px, py = x, y
            else:
                px = [self._x_viewport(xx) for xx in x]
                py = [self._y_viewport(yy) for yy in y]
            Polyline(self.slide, px, py, lcol, lthk, lpat, fcol, closed=True)
    
        else:
            xp, yp = x, y
                
            if viewport:
                pat_seg = patline.get_pattern_line(self, xp, yp, lpat.pat_len, lpat.pat_t, viewport=True)
            else:
                pat_seg = patline.get_pattern_line(self, xp, yp, lpat.pat_len, lpat.pat_t)

            for p1 in pat_seg:
                x2 = [p2[0] for p2 in p1 ]
                y2 = [p2[1] for p2 in p1 ]
             
            Polyline(self.slide, x2, y2, lcol, lthk, lpat, fcol, closed=)

def lpolygon(self, x, y, lcol=None, lthk=None, fcol=None, lpat=linepat._PAT_SOLID, viewport=False):
        pass
        
    def begin_symbol(self, sym):
        pass
        
    def end_symbol(self):
        pass        
        
    def symbol(self, x,y,sym,draw=False):
        #cx = self._x_viewport(x)
        #cy = self._y_viewport(y)
        #px, py = sym.update_xy(cx,cy)
        #ppx = [px1 for px1 in px]
        #ppy = [py1 for py1 in py]
        #self.polyline(ppx, ppy, sym.lcol, sym.lthk*self.frm.hgt(), linepat._PAT_SOLID, sym.fcol)
        pass
    
